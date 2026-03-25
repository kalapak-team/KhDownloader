import asyncio
import base64
import io

from fastapi import APIRouter, File, HTTPException, UploadFile
from fastapi.responses import JSONResponse

router = APIRouter(prefix="/api/pdf2pptx", tags=["pdf2pptx"])

MAX_PDF_SIZE_MB = 50
MAX_PDF_SIZE = MAX_PDF_SIZE_MB * 1024 * 1024


@router.post("")
async def pdf_to_pptx(
    file: UploadFile = File(..., description="PDF file to convert to PPTX"),
) -> JSONResponse:
    if not file.filename or not file.filename.lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Only PDF files are accepted")

    content = await file.read()
    if len(content) > MAX_PDF_SIZE:
        raise HTTPException(
            status_code=413,
            detail=f"PDF exceeds maximum allowed size of {MAX_PDF_SIZE_MB} MB",
        )

    def _convert() -> bytes:
        from pdf2image import convert_from_bytes  # noqa: PLC0415
        from pptx import Presentation  # noqa: PLC0415
        from pptx.util import Emu  # noqa: PLC0415

        images = convert_from_bytes(content, dpi=200, fmt="png")

        prs = Presentation()
        for img in images:
            w_px, h_px = img.size
            # Use 96 DPI reference for EMU conversion (1 inch = 914400 EMU)
            emu_w = int(w_px / 200 * 914400)
            emu_h = int(h_px / 200 * 914400)
            prs.slide_width = Emu(emu_w)
            prs.slide_height = Emu(emu_h)

            slide_layout = prs.slide_layouts[6]  # blank layout
            slide = prs.slides.add_slide(slide_layout)

            img_stream = io.BytesIO()
            img.save(img_stream, format="PNG")
            img_stream.seek(0)

            slide.shapes.add_picture(img_stream, Emu(0), Emu(0), Emu(emu_w), Emu(emu_h))

        out = io.BytesIO()
        prs.save(out)
        out.seek(0)
        return out.read()

    try:
        pptx_data = await asyncio.to_thread(_convert)
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"Conversion failed: {exc}") from exc

    stem = file.filename[:-4] if file.filename and len(file.filename) > 4 else "converted"
    safe_stem = "".join(c if ord(c) < 128 else "_" for c in stem).strip("_") or "converted"
    pptx_filename = f"{safe_stem}.pptx"

    return JSONResponse({
        "filename": pptx_filename,
        "data": base64.b64encode(pptx_data).decode("ascii"),
    })
